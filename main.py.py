import sys
from cx_Freeze import setup, Executable
from tkinter import *
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import platform
userhome = os.path.expanduser('~')



filename1='Desktop'##Masaüstü Varsayılan Olarak Ayarlandı
filename =(userhome+'\\Desktop\\PPH-10.pptx')



prs = Presentation()##Sunum Oluşturma
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)


    
def getfolder():##Dosya Yolu Bulma
    global filename,filename1
    filename = filedialog.askdirectory()+ '\\PPH-10.pptx'
    u=(len(filename))
    if    u>1:
        if len(filename)>25:
            a=(len(filename1)-len(filename1)-20)
            filename1='.../'+filename[a:]
            p.configure(text=filename1)
        else:
            p.configure(text=filename)
    else:
        filename =(userhome+'\\Desktop\\PPH-10.pptx')
        p.configure(text='Desktop')
        
    print(filename)
    print(u)

def kayit():
    global a2, metin,tf1

    ##Başlık
    baslik.text=w.get()
    baslik.alignment = PP_ALIGN.CENTER
    baslik.font.bold = True
    baslik.font.size = Pt(28)

    ##Metin
    a2=f.get(1.0,END)
    metin.text=a2
    tf1.fit_text()
    metin.font.size = Pt(18)

    #Kayıt et
    prs.save(filename)
    print(a2)
    baslik.text=''
    

##Uygulama Özellikleri
a = Tk()
a.title('PowerPoint Sunum Oluşturucu 1.0')
a.geometry('640x480')
a.resizable(width=False, height=False)

b = Frame(a,pady=5)
b.pack()

b2= Frame(a)
b2.pack()

##Başlık-Kayıt Yeri Seçimi
w = Entry(b,width=75)
w.pack(side=LEFT ,fill=X,padx=5)
s = Button(b,width=20,text='Kayıt Yeri',command=getfolder)
s.pack(side=RIGHT)

##Metin-Liste
f = Text(b2,width=56)
f.pack(side=LEFT,padx=5,fill=X)
p = Label(b2,text='Desktop',width=20)
p.pack()
l = Listbox(b2)
l.pack(fill=X)
g = Button(b2,width=20,height=5,text='Tamam',bg='lime green',command=kayit)
g.pack(pady=10)


w.insert(0,"Başlığı buraya giriniz.")
f.insert(END, "Metni buraya giriniz.")



## Başlık Oluştur
left =Inches(0)
top = Inches(0.5)
width = Inches(10)
height = Inches(0.75)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
baslik = tf.add_paragraph()

## Metin Oluştur
left1 =Inches(1)
top1 = Inches(1.5)
width1 = Inches(8)
height1 = Inches(5)
txBox1 = slide.shapes.add_textbox(left1, top1, width1, height1)
tf1 = txBox1.text_frame
metin = tf1.add_paragraph()
a2=f.get(1.0,END)
metin.text=a2











a.mainloop()


# Dependencies are automatically detected, but it might need fine tuning.
build_exe_options = {"packages": ["os"], "excludes": ["tkinter"]}

# GUI applications require a different base on Windows (the default is for a
# console application).
base = None
if sys.platform == "win32":
    base = "Win32GUI"

setup(  name = "guifoo",
        version = "0.1",
        description = "My GUI application!",
        options = {"build_exe": build_exe_options},
        executables = [Executable("guifoo.py", base=base)])
